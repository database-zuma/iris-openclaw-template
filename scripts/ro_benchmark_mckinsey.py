#!/usr/bin/env python3
"""
RO Request Benchmark — McKinsey Consulting Style Presentation
Target: Pak Steven Eka Halim (CEO Zuma Indonesia)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette (McKinsey) ──
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)
MED_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
VERY_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ZUMA_GREEN = RGBColor(0x00, 0xE2, 0x73)
RED_ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
LIGHT_BLUE = RGBColor(0xEB, 0xF5, 0xFB)
GREEN_BG = RGBColor(0xE8, 0xF8, 0xF0)

FONT_NAME = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.8)
CONTENT_W = SLIDE_W - 2 * MARGIN


# ── Helper Functions ──


def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=12,
    bold=False,
    color=DARK_GRAY,
    alignment=PP_ALIGN.LEFT,
    font_name=FONT_NAME,
    anchor=MSO_ANCHOR.TOP,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rich_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return tf


def add_para(
    tf,
    text,
    font_size=12,
    bold=False,
    color=DARK_GRAY,
    alignment=PP_ALIGN.LEFT,
    space_after=Pt(6),
    font_name=FONT_NAME,
    space_before=Pt(0),
):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    return p


def add_line(slide, left, top, width, color=NAVY, thickness=2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_page_number(slide, num, total=9):
    add_textbox(
        slide,
        SLIDE_W - Inches(1.5),
        SLIDE_H - Inches(0.5),
        Inches(1.2),
        Inches(0.3),
        f"{num}",
        font_size=9,
        color=MED_GRAY,
        alignment=PP_ALIGN.RIGHT,
    )


def add_footer_bar(slide, num, total=9):
    """Thin navy bar at bottom with page number."""
    add_line(slide, Inches(0), SLIDE_H - Inches(0.35), SLIDE_W, color=NAVY, thickness=2)
    add_textbox(
        slide,
        SLIDE_W - Inches(1.5),
        SLIDE_H - Inches(0.45),
        Inches(1.2),
        Inches(0.3),
        f"{num} / {total}",
        font_size=8,
        color=MED_GRAY,
        alignment=PP_ALIGN.RIGHT,
    )
    add_textbox(
        slide,
        MARGIN,
        SLIDE_H - Inches(0.45),
        Inches(3),
        Inches(0.3),
        "CONFIDENTIAL",
        font_size=8,
        color=MED_GRAY,
        alignment=PP_ALIGN.LEFT,
    )


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide (McKinsey Style)
# ════════════════════════════════════════════════════════════════
slide1 = add_blank_slide()

# Navy accent line at top
add_line(slide1, Inches(0), Inches(0), SLIDE_W, color=NAVY, thickness=6)

# Title — top left area
add_textbox(
    slide1,
    MARGIN,
    Inches(1.8),
    Inches(9),
    Inches(1.0),
    "RO Request Generation\nBenchmark",
    font_size=36,
    bold=True,
    color=NAVY,
)

# Subtitle
add_textbox(
    slide1,
    MARGIN,
    Inches(3.3),
    Inches(8),
    Inches(0.5),
    "Comparative Analysis — Tunjungan Plaza",
    font_size=18,
    color=MED_GRAY,
)

# Horizontal separator
add_line(slide1, MARGIN, Inches(4.0), Inches(4), color=NAVY, thickness=2)

# Bottom-right info block
info_tf = add_rich_textbox(slide1, Inches(9.0), Inches(5.0), Inches(3.8), Inches(2.0))
add_para(info_tf, "Presented to", font_size=10, color=MED_GRAY, space_after=Pt(2))
add_para(
    info_tf,
    "Pak Steven Eka Halim",
    font_size=14,
    bold=True,
    color=NAVY,
    space_after=Pt(16),
)
add_para(info_tf, "February 2026", font_size=11, color=DARK_GRAY, space_after=Pt(4))
add_para(info_tf, "Prepared by: CI Team", font_size=11, color=MED_GRAY)

# Bottom bar
add_line(slide1, Inches(0), SLIDE_H - Inches(0.35), SLIDE_W, color=NAVY, thickness=3)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Executive Summary (McKinsey Pyramid)
# ════════════════════════════════════════════════════════════════
slide2 = add_blank_slide()
add_footer_bar(slide2, 2)

# Header
add_textbox(
    slide2,
    MARGIN,
    Inches(0.4),
    Inches(8),
    Inches(0.5),
    "Executive Summary",
    font_size=24,
    bold=True,
    color=NAVY,
)
add_line(slide2, MARGIN, Inches(0.95), Inches(2.5), color=ZUMA_GREEN, thickness=3)

# Key Message Box (navy bg)
msg_box = add_rect(slide2, MARGIN, Inches(1.3), CONTENT_W, Inches(0.8), NAVY)
msg_tf = msg_box.text_frame
msg_tf.word_wrap = True
try:
    msg_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except:
    pass
p = msg_tf.paragraphs[0]
p.text = "Terminal automation is the only viable method for RO Request generation"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = FONT_NAME
p.alignment = PP_ALIGN.CENTER

# 3-column comparison
col_w = Inches(3.6)
col_gap = Inches(0.3)
col_top = Inches(2.5)
col_h = Inches(3.0)

methods = [
    (
        "Terminal (Claude CLI)",
        "✅ SUCCESS",
        ZUMA_GREEN,
        "2.66 seconds",
        "126 items (74 Protol + 52 Box)",
        "100% data completeness",
    ),
    (
        "Excel Extensions",
        "❌ FAILED",
        RED_ACCENT,
        "Not executable",
        "Configuration error",
        "API limitation on code execution",
    ),
    (
        "GSheet Browser",
        "❌ FAILED",
        RED_ACCENT,
        "~60 minutes",
        "Empty output (0 items)",
        "Cannot handle multi-sheet scale",
    ),
]

for i, (name, status, status_color, time_val, output_val, note) in enumerate(methods):
    col_left = MARGIN + i * (col_w + col_gap)

    # Card background
    add_rect(
        slide2, col_left, col_top, col_w, col_h, VERY_LIGHT, border_color=LIGHT_GRAY
    )

    # Method name
    add_textbox(
        slide2,
        col_left + Inches(0.2),
        col_top + Inches(0.15),
        col_w - Inches(0.4),
        Inches(0.35),
        name,
        font_size=13,
        bold=True,
        color=NAVY,
    )

    # Status badge
    badge = add_rect(
        slide2,
        col_left + Inches(0.2),
        col_top + Inches(0.55),
        Inches(1.6),
        Inches(0.35),
        status_color,
    )
    badge_tf = badge.text_frame
    try:
        badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    bp = badge_tf.paragraphs[0]
    bp.text = status
    bp.font.size = Pt(11)
    bp.font.bold = True
    bp.font.color.rgb = WHITE
    bp.font.name = FONT_NAME
    bp.alignment = PP_ALIGN.CENTER

    # Metrics
    metrics_tf = add_rich_textbox(
        slide2,
        col_left + Inches(0.2),
        col_top + Inches(1.1),
        col_w - Inches(0.4),
        Inches(1.8),
    )
    add_para(metrics_tf, "Time", font_size=9, color=MED_GRAY, space_after=Pt(1))
    add_para(
        metrics_tf,
        time_val,
        font_size=12,
        bold=True,
        color=DARK_GRAY,
        space_after=Pt(10),
    )
    add_para(metrics_tf, "Output", font_size=9, color=MED_GRAY, space_after=Pt(1))
    add_para(
        metrics_tf,
        output_val,
        font_size=12,
        bold=True,
        color=DARK_GRAY,
        space_after=Pt(10),
    )
    add_para(metrics_tf, note, font_size=10, color=MED_GRAY, space_after=Pt(0))

# Recommendation bar (green bg)
rec_box = add_rect(slide2, MARGIN, Inches(5.8), CONTENT_W, Inches(0.65), GREEN_BG)
rec_box.line.color.rgb = ZUMA_GREEN
rec_box.line.width = Pt(1.5)
rec_tf = rec_box.text_frame
rec_tf.word_wrap = True
try:
    rec_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except:
    pass
rp = rec_tf.paragraphs[0]
rp.text = "↗  RECOMMENDATION: Standardize Terminal approach for 60+ stores nationwide"
rp.font.size = Pt(13)
rp.font.bold = True
rp.font.color.rgb = NAVY
rp.font.name = FONT_NAME
rp.alignment = PP_ALIGN.LEFT


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Situation & Objective (SCQA)
# ════════════════════════════════════════════════════════════════
slide3 = add_blank_slide()
add_footer_bar(slide3, 3)

add_textbox(
    slide3,
    MARGIN,
    Inches(0.4),
    Inches(8),
    Inches(0.5),
    "Situation & Objective",
    font_size=24,
    bold=True,
    color=NAVY,
)
add_line(slide3, MARGIN, Inches(0.95), Inches(2.5), color=ZUMA_GREEN, thickness=3)

# Left column — SCQA (40%)
left_w = Inches(4.8)
left_x = MARGIN
scqa_top = Inches(1.4)

scqa_items = [
    (
        "SITUATION",
        "Weekly Replenishment Orders (RO) are required\nfor 60+ retail stores across Indonesia",
    ),
    (
        "COMPLICATION",
        "Manual RO generation takes ~2 hours per store,\ncreating a 120-hour weekly bottleneck",
    ),
    (
        "QUESTION",
        "Which AI-assisted method can reliably automate\nRO Request generation at scale?",
    ),
]

for i, (label, desc) in enumerate(scqa_items):
    y = scqa_top + i * Inches(1.5)

    # Label badge
    lbl = add_rect(slide3, left_x, y, Inches(1.6), Inches(0.35), NAVY)
    lbl_tf = lbl.text_frame
    try:
        lbl_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    lp = lbl_tf.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(10)
    lp.font.bold = True
    lp.font.color.rgb = WHITE
    lp.font.name = FONT_NAME
    lp.alignment = PP_ALIGN.CENTER

    # Description
    add_textbox(
        slide3,
        left_x,
        y + Inches(0.45),
        left_w,
        Inches(0.9),
        desc,
        font_size=12,
        color=DARK_GRAY,
    )

# Vertical separator
add_rect(slide3, Inches(5.9), Inches(1.4), Pt(2), Inches(4.5), LIGHT_GRAY)

# Right column — Test Setup (60%)
right_x = Inches(6.3)
right_w = Inches(6.2)

add_textbox(
    slide3,
    right_x,
    Inches(1.4),
    right_w,
    Inches(0.4),
    "TEST CASE DESIGN",
    font_size=14,
    bold=True,
    color=NAVY,
)

test_items = [
    ("Store", "Tunjungan Plaza (TP) — Jatim Branch"),
    ("Data Size", "5.6 MB across 3 input files"),
    ("Input Files", "Planogram XLSX, Stock data, Sales history"),
    ("Methods Tested", "3 — Terminal CLI, Excel Extensions, GSheet Browser"),
]

for i, (label, value) in enumerate(test_items):
    y = Inches(2.1) + i * Inches(0.65)
    # Light bg row
    if i % 2 == 0:
        add_rect(slide3, right_x, y, right_w, Inches(0.55), VERY_LIGHT)

    add_textbox(
        slide3,
        right_x + Inches(0.15),
        y + Inches(0.08),
        Inches(1.8),
        Inches(0.4),
        label,
        font_size=11,
        bold=True,
        color=NAVY,
    )
    add_textbox(
        slide3,
        right_x + Inches(2.0),
        y + Inches(0.08),
        Inches(4.0),
        Inches(0.4),
        value,
        font_size=11,
        color=DARK_GRAY,
    )

# Evaluation criteria box
eval_box = add_rect(slide3, right_x, Inches(4.8), right_w, Inches(1.0), LIGHT_BLUE)
eval_box.line.fill.background()
eval_tf = eval_box.text_frame
eval_tf.word_wrap = True
add_para(
    eval_tf,
    "EVALUATION CRITERIA",
    font_size=10,
    bold=True,
    color=NAVY,
    space_after=Pt(6),
)
add_para(
    eval_tf,
    "Speed  |  Output Quality  |  Reliability  |  Scalability",
    font_size=12,
    color=DARK_GRAY,
)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Results Overview (Table)
# ════════════════════════════════════════════════════════════════
slide4 = add_blank_slide()
add_footer_bar(slide4, 4)

add_textbox(
    slide4,
    MARGIN,
    Inches(0.4),
    Inches(8),
    Inches(0.5),
    "Results Overview",
    font_size=24,
    bold=True,
    color=NAVY,
)
add_line(slide4, MARGIN, Inches(0.95), Inches(2.5), color=ZUMA_GREEN, thickness=3)

# Table
table_top = Inches(1.4)
table_left = MARGIN
table_w = CONTENT_W
row_h = Inches(0.55)
cols = ["Method", "Status", "Time", "Output", "Quality Score"]
col_widths = [Inches(2.8), Inches(1.5), Inches(2.0), Inches(3.0), Inches(2.4)]

# Header row
header_bg = add_rect(slide4, table_left, table_top, table_w, row_h, NAVY)
x_offset = table_left
for j, (col_name, cw) in enumerate(zip(cols, col_widths)):
    add_textbox(
        slide4,
        x_offset + Inches(0.15),
        table_top + Inches(0.08),
        cw - Inches(0.3),
        row_h - Inches(0.15),
        col_name,
        font_size=11,
        bold=True,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    x_offset += cw

# Data rows
rows_data = [
    (
        "Terminal (Claude CLI)",
        "✅ Success",
        "2.66 seconds",
        "74 Protol + 52 Box",
        "100%",
    ),
    ("Excel Extensions", "❌ Failed", "—", "Config Error", "N/A"),
    ("GSheet Browser", "❌ Failed", "~60 minutes", "Empty (TBD)", "N/A"),
]

for i, row in enumerate(rows_data):
    y = table_top + (i + 1) * row_h
    bg_color = WHITE if i % 2 == 0 else VERY_LIGHT
    add_rect(slide4, table_left, y, table_w, row_h, bg_color, border_color=LIGHT_GRAY)

    x_offset = table_left
    for j, (val, cw) in enumerate(zip(row, col_widths)):
        txt_color = DARK_GRAY
        txt_bold = False
        if j == 0:
            txt_bold = True
            txt_color = NAVY
        elif j == 1:
            if "✅" in val:
                txt_color = ZUMA_GREEN
            else:
                txt_color = RED_ACCENT
            txt_bold = True
        add_textbox(
            slide4,
            x_offset + Inches(0.15),
            y + Inches(0.08),
            cw - Inches(0.3),
            row_h - Inches(0.15),
            val,
            font_size=11,
            bold=txt_bold,
            color=txt_color,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        x_offset += cw

# Bottom border
add_line(slide4, table_left, table_top + 4 * row_h, table_w, color=NAVY, thickness=2)

# Key Insight callout
insight_y = table_top + 4 * row_h + Inches(0.5)
insight_box = add_rect(slide4, table_left, insight_y, table_w, Inches(0.8), LIGHT_BLUE)
insight_box.line.fill.background()
insight_tf = insight_box.text_frame
insight_tf.word_wrap = True
try:
    insight_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except:
    pass
add_para(
    insight_tf, "KEY INSIGHT", font_size=10, bold=True, color=NAVY, space_after=Pt(4)
)
add_para(
    insight_tf,
    "Terminal is 1,353× faster than GSheet with 100% data completeness — the only method that produced usable output",
    font_size=13,
    color=DARK_GRAY,
)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Terminal Method Deep Dive
# ════════════════════════════════════════════════════════════════
slide5 = add_blank_slide()
add_footer_bar(slide5, 5)

add_textbox(
    slide5,
    MARGIN,
    Inches(0.4),
    Inches(8),
    Inches(0.5),
    "Terminal Automation — Deep Dive",
    font_size=24,
    bold=True,
    color=NAVY,
)

# Success badge
badge5 = add_rect(
    slide5, Inches(8.5), Inches(0.4), Inches(1.6), Inches(0.4), ZUMA_GREEN
)
badge5_tf = badge5.text_frame
try:
    badge5_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except:
    pass
bp5 = badge5_tf.paragraphs[0]
bp5.text = "✅ SUCCESS"
bp5.font.size = Pt(12)
bp5.font.bold = True
bp5.font.color.rgb = WHITE
bp5.font.name = FONT_NAME
bp5.alignment = PP_ALIGN.CENTER

add_line(slide5, MARGIN, Inches(0.95), Inches(2.5), color=ZUMA_GREEN, thickness=3)

# 4 Metric Cards
card_w = Inches(2.7)
card_h = Inches(2.0)
card_gap = Inches(0.3)
cards_top = Inches(1.4)

metrics = [
    ("SPEED", "2.66s", "Total execution time", NAVY),
    ("OUTPUT", "126", "Items generated (74P + 52B)", NAVY),
    ("MODEL", "Kimi K2.5", "Claude CLI engine", NAVY),
    ("ACCURACY", "100%", "Data completeness", ZUMA_GREEN),
]

for i, (label, value, sublabel, accent) in enumerate(metrics):
    cx = MARGIN + i * (card_w + card_gap)

    # Card bg
    card = add_rect(
        slide5, cx, cards_top, card_w, card_h, VERY_LIGHT, border_color=LIGHT_GRAY
    )

    # Top accent line
    add_rect(slide5, cx, cards_top, card_w, Pt(4), accent)

    # Label
    add_textbox(
        slide5,
        cx + Inches(0.2),
        cards_top + Inches(0.25),
        card_w - Inches(0.4),
        Inches(0.3),
        label,
        font_size=10,
        bold=True,
        color=MED_GRAY,
    )

    # Big number
    add_textbox(
        slide5,
        cx + Inches(0.2),
        cards_top + Inches(0.6),
        card_w - Inches(0.4),
        Inches(0.7),
        value,
        font_size=32,
        bold=True,
        color=accent,
    )

    # Sub-label
    add_textbox(
        slide5,
        cx + Inches(0.2),
        cards_top + Inches(1.4),
        card_w - Inches(0.4),
        Inches(0.4),
        sublabel,
        font_size=10,
        color=MED_GRAY,
    )

# Key Success Factors section
ksf_top = Inches(3.8)
add_textbox(
    slide5,
    MARGIN,
    ksf_top,
    Inches(4),
    Inches(0.4),
    "KEY SUCCESS FACTORS",
    font_size=14,
    bold=True,
    color=NAVY,
)
add_line(
    slide5, MARGIN, ksf_top + Inches(0.4), Inches(1.5), color=ZUMA_GREEN, thickness=2
)

factors = [
    (
        "Direct Database Access",
        "Queries PostgreSQL directly — no file upload/download bottleneck",
    ),
    (
        "Programmatic Logic",
        "Python-based RO calculation — no UI or spreadsheet constraints",
    ),
    ("Reliable & Repeatable", "Same script, same output every time — zero human error"),
]

for i, (title, desc) in enumerate(factors):
    fy = ksf_top + Inches(0.7) + i * Inches(0.75)
    # Navy bullet dot
    add_rect(
        slide5,
        MARGIN + Inches(0.1),
        fy + Inches(0.08),
        Inches(0.12),
        Inches(0.12),
        NAVY,
    )
    add_textbox(
        slide5,
        MARGIN + Inches(0.4),
        fy - Inches(0.02),
        Inches(4),
        Inches(0.3),
        title,
        font_size=12,
        bold=True,
        color=NAVY,
    )
    add_textbox(
        slide5,
        MARGIN + Inches(0.4),
        fy + Inches(0.25),
        Inches(10),
        Inches(0.3),
        desc,
        font_size=11,
        color=DARK_GRAY,
    )


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Excel Extensions Analysis
# ════════════════════════════════════════════════════════════════
slide6 = add_blank_slide()
add_footer_bar(slide6, 6)

add_textbox(
    slide6,
    MARGIN,
    Inches(0.4),
    Inches(8),
    Inches(0.5),
    "Excel Extensions — Analysis",
    font_size=24,
    bold=True,
    color=NAVY,
)

# Failed badge
badge6 = add_rect(
    slide6, Inches(8.0), Inches(0.4), Inches(1.4), Inches(0.4), RED_ACCENT
)
badge6_tf = badge6.text_frame
try:
    badge6_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except:
    pass
bp6 = badge6_tf.paragraphs[0]
bp6.text = "❌ FAILED"
bp6.font.size = Pt(12)
bp6.font.bold = True
bp6.font.color.rgb = WHITE
bp6.font.name = FONT_NAME
bp6.alignment = PP_ALIGN.CENTER

add_line(slide6, MARGIN, Inches(0.95), Inches(2.5), color=RED_ACCENT, thickness=3)

# Error message box
err_top = Inches(1.5)
err_box = add_rect(
    slide6,
    MARGIN,
    err_top,
    CONTENT_W,
    Inches(1.2),
    RGBColor(0xF5, 0xF5, 0xF5),
    border_color=LIGHT_GRAY,
)
err_tf = err_box.text_frame
err_tf.word_wrap = True
add_para(
    err_tf, "ERROR OUTPUT", font_size=9, bold=True, color=RED_ACCENT, space_after=Pt(8)
)
add_para(
    err_tf,
    "configurations.1.code_execution: Extra inputs are not permitted",
    font_size=14,
    color=DARK_GRAY,
    font_name="Courier New",
)

# Analysis section
analysis_top = Inches(3.1)
sections = [
    (
        "ROOT CAUSE",
        "The ChatGPT API does not support 'code_execution' as a configuration\nparameter for Excel file processing. This is an API schema limitation\nthat cannot be resolved through prompt engineering.",
    ),
    (
        "WHAT WAS TESTED",
        "Multiple ChatGPT models (GPT-4o, GPT-4 Turbo) with Excel file uploads\nand various configuration schemas for code execution.",
    ),
    (
        "CONCLUSION",
        "Excel-based AI extensions are fundamentally incompatible with\nthe multi-sheet calculation requirements of RO generation.",
    ),
]

for i, (label, desc) in enumerate(sections):
    sy = analysis_top + i * Inches(1.25)

    # Left label
    lbl = add_rect(slide6, MARGIN, sy, Inches(1.8), Inches(0.32), NAVY)
    lbl_tf = lbl.text_frame
    try:
        lbl_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    lp = lbl_tf.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(9)
    lp.font.bold = True
    lp.font.color.rgb = WHITE
    lp.font.name = FONT_NAME
    lp.alignment = PP_ALIGN.CENTER

    # Description
    add_textbox(
        slide6,
        MARGIN,
        sy + Inches(0.42),
        CONTENT_W,
        Inches(0.75),
        desc,
        font_size=11,
        color=DARK_GRAY,
    )


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — GSheet Browser Analysis
# ════════════════════════════════════════════════════════════════
slide7 = add_blank_slide()
add_footer_bar(slide7, 7)

add_textbox(
    slide7,
    MARGIN,
    Inches(0.4),
    Inches(8),
    Inches(0.5),
    "GSheet Browser Extensions — Analysis",
    font_size=24,
    bold=True,
    color=NAVY,
)

badge7 = add_rect(
    slide7, Inches(8.8), Inches(0.4), Inches(1.4), Inches(0.4), RED_ACCENT
)
badge7_tf = badge7.text_frame
try:
    badge7_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except:
    pass
bp7 = badge7_tf.paragraphs[0]
bp7.text = "❌ FAILED"
bp7.font.size = Pt(12)
bp7.font.bold = True
bp7.font.color.rgb = WHITE
bp7.font.name = FONT_NAME
bp7.alignment = PP_ALIGN.CENTER

add_line(slide7, MARGIN, Inches(0.95), Inches(2.5), color=RED_ACCENT, thickness=3)

# Left column — Performance
left7_x = MARGIN
left7_w = Inches(5.5)
perf_top = Inches(1.4)

add_textbox(
    slide7,
    left7_x,
    perf_top,
    Inches(3),
    Inches(0.35),
    "PERFORMANCE METRICS",
    font_size=12,
    bold=True,
    color=NAVY,
)

perf_metrics = [
    ("Execution Time", "~60 minutes"),
    ("Speed Ratio", "1,353× slower than Terminal"),
    ("Token Usage", "High (browser automation overhead)"),
]
for i, (label, val) in enumerate(perf_metrics):
    py = perf_top + Inches(0.6) + i * Inches(0.6)
    bg_c = VERY_LIGHT if i % 2 == 0 else WHITE
    add_rect(slide7, left7_x, py, left7_w, Inches(0.5), bg_c, border_color=LIGHT_GRAY)
    add_textbox(
        slide7,
        left7_x + Inches(0.15),
        py + Inches(0.05),
        Inches(2.2),
        Inches(0.4),
        label,
        font_size=11,
        bold=True,
        color=NAVY,
    )
    add_textbox(
        slide7,
        left7_x + Inches(2.4),
        py + Inches(0.05),
        Inches(3.0),
        Inches(0.4),
        val,
        font_size=11,
        color=DARK_GRAY,
    )

# Right column — Output Quality
right7_x = Inches(7.0)
right7_w = Inches(5.5)
add_textbox(
    slide7,
    right7_x,
    perf_top,
    Inches(3),
    Inches(0.35),
    "OUTPUT QUALITY",
    font_size=12,
    bold=True,
    color=NAVY,
)

output_metrics = [
    ("RO Protol", "TBD — Empty output"),
    ("RO Box", "TBD — Empty output"),
    ("Surplus Pull", "0 items — Empty output"),
]
for i, (label, val) in enumerate(output_metrics):
    py = perf_top + Inches(0.6) + i * Inches(0.6)
    bg_c = VERY_LIGHT if i % 2 == 0 else WHITE
    add_rect(slide7, right7_x, py, right7_w, Inches(0.5), bg_c, border_color=LIGHT_GRAY)
    add_textbox(
        slide7,
        right7_x + Inches(0.15),
        py + Inches(0.05),
        Inches(2.2),
        Inches(0.4),
        label,
        font_size=11,
        bold=True,
        color=NAVY,
    )
    add_textbox(
        slide7,
        right7_x + Inches(2.4),
        py + Inches(0.05),
        Inches(3.0),
        Inches(0.4),
        val,
        font_size=11,
        color=RED_ACCENT,
        bold=True,
    )

# Root cause section
rc_top = Inches(3.6)
rc_box = add_rect(
    slide7, MARGIN, rc_top, CONTENT_W, Inches(1.6), VERY_LIGHT, border_color=LIGHT_GRAY
)
rc_tf = rc_box.text_frame
rc_tf.word_wrap = True
add_para(
    rc_tf,
    "ROOT CAUSE ANALYSIS",
    font_size=12,
    bold=True,
    color=NAVY,
    space_after=Pt(10),
)
add_para(
    rc_tf,
    "1.  Browser-based AI cannot process multi-sheet Excel files natively — requires upload to GSheet first",
    font_size=11,
    color=DARK_GRAY,
    space_after=Pt(6),
)
add_para(
    rc_tf,
    "2.  GSheet formulas and cross-references lost during conversion — breaks RO calculation logic",
    font_size=11,
    color=DARK_GRAY,
    space_after=Pt(6),
)
add_para(
    rc_tf,
    "3.  Browser automation overhead (DOM parsing, UI interactions) adds ~60 min of latency per store",
    font_size=11,
    color=DARK_GRAY,
    space_after=Pt(6),
)
add_para(
    rc_tf,
    "4.  No programmatic access to database — all data must flow through browser UI",
    font_size=11,
    color=DARK_GRAY,
)

# Verdict
verdict_y = Inches(5.6)
verdict = add_rect(
    slide7, MARGIN, verdict_y, CONTENT_W, Inches(0.6), RGBColor(0xFD, 0xED, 0xED)
)
verdict.line.color.rgb = RED_ACCENT
verdict.line.width = Pt(1)
verdict_tf = verdict.text_frame
try:
    verdict_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except:
    pass
vp = verdict_tf.paragraphs[0]
vp.text = "VERDICT: Not viable at scale — 60 stores × 60 min = 60 hours/week (vs 3 min with Terminal)"
vp.font.size = Pt(12)
vp.font.bold = True
vp.font.color.rgb = RED_ACCENT
vp.font.name = FONT_NAME
vp.alignment = PP_ALIGN.LEFT


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Business Impact (McKinsey Pyramid)
# ════════════════════════════════════════════════════════════════
slide8 = add_blank_slide()
add_footer_bar(slide8, 8)

add_textbox(
    slide8,
    MARGIN,
    Inches(0.4),
    Inches(8),
    Inches(0.5),
    "Business Impact",
    font_size=24,
    bold=True,
    color=NAVY,
)
add_line(slide8, MARGIN, Inches(0.95), Inches(2.5), color=ZUMA_GREEN, thickness=3)

# Impact headline
impact_box = add_rect(slide8, MARGIN, Inches(1.3), CONTENT_W, Inches(0.7), NAVY)
impact_tf = impact_box.text_frame
try:
    impact_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except:
    pass
ip = impact_tf.paragraphs[0]
ip.text = "Terminal automation delivers 99.96% time reduction in RO Request generation"
ip.font.size = Pt(16)
ip.font.bold = True
ip.font.color.rgb = WHITE
ip.font.name = FONT_NAME
ip.alignment = PP_ALIGN.CENTER

# Two-column comparison
comp_top = Inches(2.4)
comp_w = Inches(5.5)
comp_h = Inches(2.5)

# CURRENT state
cur_box = add_rect(
    slide8, MARGIN, comp_top, comp_w, comp_h, VERY_LIGHT, border_color=LIGHT_GRAY
)
add_rect(slide8, MARGIN, comp_top, comp_w, Pt(4), RED_ACCENT)

cur_tf = add_rich_textbox(
    slide8,
    MARGIN + Inches(0.3),
    comp_top + Inches(0.3),
    comp_w - Inches(0.6),
    comp_h - Inches(0.5),
)
add_para(
    cur_tf,
    "CURRENT STATE",
    font_size=11,
    bold=True,
    color=RED_ACCENT,
    space_after=Pt(12),
)
add_para(
    cur_tf,
    "60 stores  ×  ~2 hours  =",
    font_size=14,
    color=DARK_GRAY,
    space_after=Pt(4),
)
add_para(
    cur_tf,
    "120 hours / week",
    font_size=28,
    bold=True,
    color=RED_ACCENT,
    space_after=Pt(12),
)
add_para(
    cur_tf, "Manual process, error-prone, inconsistent", font_size=11, color=MED_GRAY
)

# Arrow between columns
arrow_x = MARGIN + comp_w + Inches(0.15)
add_textbox(
    slide8,
    arrow_x,
    comp_top + Inches(0.8),
    Inches(1.0),
    Inches(0.8),
    "→",
    font_size=40,
    bold=True,
    color=NAVY,
    alignment=PP_ALIGN.CENTER,
)

# FUTURE state
fut_x = MARGIN + comp_w + Inches(1.0)
fut_box = add_rect(
    slide8, fut_x, comp_top, comp_w, comp_h, GREEN_BG, border_color=ZUMA_GREEN
)
add_rect(slide8, fut_x, comp_top, comp_w, Pt(4), ZUMA_GREEN)

fut_tf = add_rich_textbox(
    slide8,
    fut_x + Inches(0.3),
    comp_top + Inches(0.3),
    comp_w - Inches(0.6),
    comp_h - Inches(0.5),
)
add_para(
    fut_tf,
    "FUTURE STATE (Terminal)",
    font_size=11,
    bold=True,
    color=ZUMA_GREEN,
    space_after=Pt(12),
)
add_para(
    fut_tf,
    "60 stores  ×  2.66 sec  =",
    font_size=14,
    color=DARK_GRAY,
    space_after=Pt(4),
)
add_para(
    fut_tf,
    "~3 minutes / week",
    font_size=28,
    bold=True,
    color=ZUMA_GREEN,
    space_after=Pt(12),
)
add_para(fut_tf, "Automated, accurate, fully repeatable", font_size=11, color=MED_GRAY)

# Recommendations (3 boxes with green left border)
rec_top = Inches(5.3)
add_textbox(
    slide8,
    MARGIN,
    rec_top,
    Inches(4),
    Inches(0.35),
    "RECOMMENDATIONS",
    font_size=12,
    bold=True,
    color=NAVY,
)

recs = [
    ("1", "Standardize Terminal approach as the official RO generation method"),
    ("2", "Scale automation to all 60+ stores with weekly cron scheduling"),
    ("3", "Retire Excel and GSheet methods — redirect resources to Terminal pipeline"),
]

for i, (num, text) in enumerate(recs):
    ry = rec_top + Inches(0.5) + i * Inches(0.55)
    # Green left border
    add_rect(slide8, MARGIN, ry, Inches(0.06), Inches(0.45), ZUMA_GREEN)
    # Number
    add_textbox(
        slide8,
        MARGIN + Inches(0.2),
        ry + Inches(0.03),
        Inches(0.3),
        Inches(0.4),
        num + ".",
        font_size=12,
        bold=True,
        color=NAVY,
    )
    # Text
    add_textbox(
        slide8,
        MARGIN + Inches(0.5),
        ry + Inches(0.03),
        Inches(11),
        Inches(0.4),
        text,
        font_size=12,
        color=DARK_GRAY,
    )


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Implementation Roadmap
# ════════════════════════════════════════════════════════════════
slide9 = add_blank_slide()
add_footer_bar(slide9, 9)

add_textbox(
    slide9,
    MARGIN,
    Inches(0.4),
    Inches(8),
    Inches(0.5),
    "Implementation Roadmap",
    font_size=24,
    bold=True,
    color=NAVY,
)
add_line(slide9, MARGIN, Inches(0.95), Inches(2.5), color=ZUMA_GREEN, thickness=3)

# Timeline
phase_w = Inches(3.6)
phase_h = Inches(4.2)
phase_gap = Inches(0.35)
phase_top = Inches(1.4)

phases = [
    (
        "IMMEDIATE",
        "This Week",
        ZUMA_GREEN,
        [
            "Finalize Terminal script for all stores",
            "Validate output format with OPS team",
            "Train OPS team on script execution",
            "Document process & error handling",
        ],
    ),
    (
        "SHORT-TERM",
        "This Month",
        NAVY,
        [
            "Deploy weekly cron automation",
            "Set up monitoring & alerting",
            "Build performance dashboard",
            "Pilot with 5 high-volume stores",
        ],
    ),
    (
        "LONG-TERM",
        "Q1 2026",
        MED_GRAY,
        [
            "Roll out to all 60+ stores",
            "Expand to Planogram generation",
            "Expand to Surplus Pull automation",
            "Integrate with Accurate Online ERP",
        ],
    ),
]

# Horizontal connecting line
line_y = phase_top + Inches(0.45)
add_rect(slide9, MARGIN + Inches(0.4), line_y, Inches(11.2), Pt(3), LIGHT_GRAY)

for i, (title, subtitle, accent, items) in enumerate(phases):
    px = MARGIN + i * (phase_w + phase_gap)

    # Phase dot on timeline
    dot_size = Inches(0.25)
    dot = add_rect(
        slide9, px + Inches(0.28), line_y - Inches(0.07), dot_size, dot_size, accent
    )

    # Phase label
    phase_label_y = phase_top + Inches(0.9)
    lbl = add_rect(slide9, px, phase_label_y, Inches(2.0), Inches(0.35), accent)
    lbl_tf = lbl.text_frame
    try:
        lbl_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    lp = lbl_tf.paragraphs[0]
    lp.text = title
    lp.font.size = Pt(10)
    lp.font.bold = True
    lp.font.color.rgb = WHITE
    lp.font.name = FONT_NAME
    lp.alignment = PP_ALIGN.CENTER

    # Subtitle
    add_textbox(
        slide9,
        px + Inches(2.1),
        phase_label_y,
        Inches(1.4),
        Inches(0.35),
        subtitle,
        font_size=10,
        color=MED_GRAY,
        alignment=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Card
    card_y = phase_label_y + Inches(0.55)
    card = add_rect(
        slide9, px, card_y, phase_w, Inches(2.5), VERY_LIGHT, border_color=LIGHT_GRAY
    )
    add_rect(slide9, px, card_y, phase_w, Pt(3), accent)

    # Items
    items_tf = add_rich_textbox(
        slide9,
        px + Inches(0.2),
        card_y + Inches(0.2),
        phase_w - Inches(0.4),
        Inches(2.2),
    )
    for item in items:
        add_para(
            items_tf, f"▸  {item}", font_size=11, color=DARK_GRAY, space_after=Pt(8)
        )

# Bottom note
note_y = Inches(6.3)
note_box = add_rect(slide9, MARGIN, note_y, CONTENT_W, Inches(0.55), LIGHT_BLUE)
note_box.line.fill.background()
note_tf = note_box.text_frame
try:
    note_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
except:
    pass
np = note_tf.paragraphs[0]
np.text = "Target: Full automation operational by end of Q1 2026 — projected savings of 6,000+ man-hours annually"
np.font.size = Pt(12)
np.font.bold = True
np.font.color.rgb = NAVY
np.font.name = FONT_NAME
np.alignment = PP_ALIGN.CENTER


# ── Save ──
output_path = os.path.expanduser("~/Desktop/RO_Request_Benchmark_McKinsey.pptx")
# Fallback if Desktop not accessible
try:
    prs.save(output_path)
    print(f"Saved: {output_path}")
except PermissionError:
    output_path = (
        "/Users/database-zuma/.openclaw/workspace/RO_Request_Benchmark_McKinsey.pptx"
    )
    prs.save(output_path)
    print(f"Saved (fallback): {output_path}")

print("Done — 9 slides generated")
